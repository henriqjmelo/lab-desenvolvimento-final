#!/usr/bin/env python3
"""
Small utility to generate a PowerPoint presentation for the RU - Restaurante Universitário project.
It will try to convert the SVG logo to PNG (using cairosvg) and embed it. If conversion is not available,
place a PNG logo at `imagens/logo.png` before running.

Usage:
  python tools/generate_presentation.py

Output: presentation/ru_presentation.pptx
"""

from pathlib import Path
import sys

BASE = Path(__file__).resolve().parent.parent
LOGO_SVG = BASE / 'src' / 'main' / 'resources' / 'static' / 'images' / 'logo.svg'
LOGO_PNG = BASE / 'imagens' / 'logo.png'
OUT_DIR = BASE / 'presentation'
OUT_DIR.mkdir(exist_ok=True)
OUT_FILE = OUT_DIR / 'ru_presentation.pptx'

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception as e:
    print('Missing dependency: python-pptx. Install with: pip install python-pptx')
    sys.exit(1)

# Try to convert svg -> png if possible
if LOGO_SVG.exists():
    try:
        import cairosvg
    except Exception:
        cairosvg = None

# Convert if cairosvg is available
if LOGO_SVG.exists() and 'cairosvg' in globals() and cairosvg is not None:
    try:
        LOGO_PNG.parent.mkdir(parents=True, exist_ok=True)
        cairosvg.svg2png(url=str(LOGO_SVG), write_to=str(LOGO_PNG))
        print('Converted SVG to PNG at', LOGO_PNG)
    except Exception as e:
        print('Failed to convert SVG to PNG:', e)

# Ensure we have a PNG to embed
if not LOGO_PNG.exists():
    print('\nLogo PNG not found at', LOGO_PNG)
    print('Please provide a PNG logo at imagens/logo.png or install cairosvg to auto-convert the SVG.')
    # continue and create a presentation without logo

prs = Presentation()

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
if LOGO_PNG.exists():
    left = Inches(0.5)
    top = Inches(0.5)
    height = Inches(1)
    slide.shapes.add_picture(str(LOGO_PNG), left, top, height=height)

title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = 'RU — Restaurante Universitário'
subtitle.text = 'Sistema de Moeda Estudantil — Apresentação'

# Overview slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = 'Visão Geral'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Objetivos do sistema:'
for point in ['Reconhecimento por mérito acadêmico', 'Envio e resgate de moedas', 'Integração com empresas parceiras']:
    p = body.add_paragraph()
    p.text = '• ' + point
    p.level = 1

# Arquitetura slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = 'Arquitetura (resumida)'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Componentes principais:'
for point in ['Módulo de Autenticação e Usuários', 'Módulo de Moedas e Transações', 'Módulo de Vantagens e Resgates', 'Banco H2 (dev)']:
    p = body.add_paragraph()
    p.text = '• ' + point
    p.level = 1

# Placeholders for screenshots
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = 'Capturas de Tela'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Adicione capturas de tela na pasta `imagens/` e re-run este script para incorporá-las.'

# Next steps slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = 'Próximos Passos'
body = slide.shapes.placeholders[1].text_frame
for point in ['Refinar UI/UX', 'Adicionar testes automatizados', 'Preparar implantação (Docker/Azure)']:
    p = body.add_paragraph()
    p.text = '• ' + point
    p.level = 1

prs.save(str(OUT_FILE))
print('\nPresentation generated:', OUT_FILE)
